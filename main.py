import os
import logging
import asyncio
import json
import uuid
from datetime import datetime, timedelta
from typing import Optional, Dict, Any, List

import aiosqlite
import google.generativeai as genai
from aiogram import Bot, Dispatcher, types, Router, F
from aiogram.filters import CommandStart, Command
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.utils.keyboard import ReplyKeyboardBuilder, InlineKeyboardBuilder
from aiogram.types import WebAppInfo, Message, CallbackQuery
import uvicorn
from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from dotenv import load_dotenv

# Load environment variables
load_dotenv('config.env', override=True)

# Configuration
BOT_TOKEN = os.getenv("BOT_TOKEN")
if BOT_TOKEN:
    print(f"DEBUG: Loaded token: '{BOT_TOKEN}' (length: {len(BOT_TOKEN)})")
else:
    print("DEBUG: BOT_TOKEN is None or Empty")
ADMIN_ID = int(os.getenv("ADMIN_ID", "0"))
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
PREMIUM_PRICE = int(os.getenv("PREMIUM_PRICE", "50000"))
WEBAPP_URL = os.getenv("WEBAPP_URL", "")

# Configure Logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Configure Gemini AI
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-pro')
else:
    logger.warning("GEMINI_API_KEY not found in environment variables!")
    model = None

# Database Manager
class Database:
    def __init__(self, db_name='talaba_bot.db'):
        self.db_name = db_name

    async def init_db(self):
        async with aiosqlite.connect(self.db_name) as db:
            # Users table
            await db.execute('''CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                tg_id INTEGER UNIQUE,
                username TEXT,
                full_name TEXT,
                is_premium INTEGER DEFAULT 0,
                premium_until TEXT,
                created TEXT,
                balance INTEGER DEFAULT 0,
                content_count INTEGER DEFAULT 0
            )''')
            
            # Payments table
            await db.execute('''CREATE TABLE IF NOT EXISTS payments (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                tg_id INTEGER,
                amount INTEGER,
                card_number TEXT,
                card_holder TEXT,
                status TEXT DEFAULT 'pending',
                payment_id TEXT UNIQUE,
                created TEXT,
                approved TEXT,
                admin_note TEXT
            )''')
            
            # Generated content table
            await db.execute('''CREATE TABLE IF NOT EXISTS generated_content (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                tg_id INTEGER,
                content_type TEXT,
                topic TEXT,
                content TEXT,
                created TEXT,
                quality TEXT DEFAULT 'standard'
            )''')
            await db.commit()

    async def get_user(self, tg_id: int):
        async with aiosqlite.connect(self.db_name) as db:
            db.row_factory = aiosqlite.Row
            cursor = await db.execute("SELECT * FROM users WHERE tg_id=?", (tg_id,))
            return await cursor.fetchone()

    async def create_or_update_user(self, tg_id: int, username: str, full_name: str):
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute('''INSERT OR REPLACE INTO users (tg_id, username, full_name, created)
                VALUES (?, ?, ?, COALESCE((SELECT created FROM users WHERE tg_id=?), ?))''',
                (tg_id, username, full_name, tg_id, datetime.now().isoformat()))
            await db.commit()

    async def update_user_balance(self, tg_id: int, amount: int):
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute("UPDATE users SET balance = balance + ? WHERE tg_id=?", (amount, tg_id))
            await db.commit()

    async def set_premium(self, tg_id: int, days: int = 30):
        premium_until = (datetime.now() + timedelta(days=days)).isoformat()
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute("UPDATE users SET is_premium=1, premium_until=? WHERE tg_id=?", (premium_until, tg_id))
            await db.commit()

    async def increment_content_count(self, tg_id: int):
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute("UPDATE users SET content_count = content_count + 1 WHERE tg_id=?", (tg_id,))
            await db.commit()

    async def create_payment(self, tg_id: int, amount: int, card_number: str, card_holder: str, payment_id: str):
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute('''INSERT INTO payments (tg_id, amount, card_number, card_holder, payment_id, created)
                VALUES (?, ?, ?, ?, ?, ?)''',
                (tg_id, amount, card_number, card_holder, payment_id, datetime.now().isoformat()))
            await db.commit()

    async def get_pending_payments(self):
        async with aiosqlite.connect(self.db_name) as db:
            db.row_factory = aiosqlite.Row
            cursor = await db.execute("SELECT * FROM payments WHERE status='pending'")
            return await cursor.fetchall()

    async def get_user_payments(self, tg_id: int):
        async with aiosqlite.connect(self.db_name) as db:
            db.row_factory = aiosqlite.Row
            cursor = await db.execute("SELECT * FROM payments WHERE tg_id=? ORDER BY created DESC", (tg_id,))
            return await cursor.fetchall()
            
    async def get_payment(self, payment_id: str):
        async with aiosqlite.connect(self.db_name) as db:
            db.row_factory = aiosqlite.Row
            cursor = await db.execute("SELECT * FROM payments WHERE payment_id=?", (payment_id,))
            return await cursor.fetchone()

    async def update_payment_status(self, payment_id: str, status: str):
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute("UPDATE payments SET status=?, approved=? WHERE payment_id=?", 
                             (status, datetime.now().isoformat(), payment_id))
            await db.commit()

    async def save_content(self, tg_id: int, content_type: str, topic: str, content: str):
        async with aiosqlite.connect(self.db_name) as db:
            await db.execute('''INSERT INTO generated_content (tg_id, content_type, topic, content, created)
                VALUES (?, ?, ?, ?, ?)''',
                (tg_id, content_type, topic, content, datetime.now().isoformat()))
            await db.commit()
            
    async def get_stats(self):
        async with aiosqlite.connect(self.db_name) as db:
            users = await (await db.execute("SELECT COUNT(*) FROM users")).fetchone()
            payments = await (await db.execute("SELECT COUNT(*) FROM payments WHERE status='pending'")).fetchone()
            content = await (await db.execute("SELECT COUNT(*) FROM generated_content")).fetchone()
            return {
                "users": users[0],
                "pending_payments": payments[0],
                "content": content[0]
            }

db = Database()

# AI Service
class AIService:
    @staticmethod
    async def generate_content(prompt: str) -> str:
        if not model:
            return "❌ AI Xizmati vaqtincha ishlamayapti (API Key topilmadi)."
        try:
            response = await model.generate_content_async(prompt)
            return response.text
        except Exception as e:
            logger.error(f"AI Generation Error: {e}")
            return "❌ Kechirasiz, kontent yaratishda xatolik yuz berdi. Iltimos keyinroq urinib ko'ring."

    @staticmethod
    async def generate_referat(topic: str, size: str) -> str:
        prompt = f"""
        Mavzu: {topic}
        Hajmi: {size} betga mo'ljallangan
        Vazifa: O'zbek tilida professional, akademik uslubda referat yozing.
        
        Tuzilishi:
        1. Kirish (Mavzuning dolzarbligi)
        2. Asosiy qism (Kamida 3 ta reja asosida)
        3. Tahliliy qism (Muammolar va yechimlar)
        4. Xulosa
        5. Foydalanilgan adabiyotlar (Kamida 5 ta manba)
        
        Formatlash: Markdown formatidan foydalaning, sarlavhalarni ajratib ko'rsating.
        """
        return await AIService.generate_content(prompt)

    @staticmethod
    async def generate_presentation(topic: str, slides: str) -> str:
        prompt = f"""
        Role: Sen professional notiq va prezentatsiya mutaxassisisan.
        Vazifa: "{topic}" mavzusida {slides} ta slayddan iborat professional prezentatsiya rejasini va matnini tayyorla.

        TALABLAR:
        1. Har bir slayd aniq ajratilgan bo'lsin.
        2. Slayd tuzilishi quyidagicha bo'lsin:
           - **Slayd Raqami va Sarlavhasi** (Markdown H2 yoki H3 bilan)
           - **Vizual Tushuncha:** (Ushbu slayd uchun qanday rasm yoki grafika mos kelishini qisqa tasvirlab ber)
           - **Asosiy Matn (Slayd uchun):** (Qisqa, lo'nda, o'qilishi oson bullet pointlar)
           - **Spiker Nutqi:** (Taqdimotchi gapirishi uchun jonli va ishonchli matn)

        Formatlash: Toza Markdown formatidan foydalan. Sarlavhalar uchun #, qalin yozuv uchun ** belgilarini ishlat.
        Til: To'liq O'zbek tilida (Kirill yoki Lotin alifbosida foydalanuvchi so'roviga moslab, hozir Lotincha yoz).
        """
        return await AIService.generate_content(prompt)
        
    @staticmethod
    async def generate_essay(topic: str, essay_type: str) -> str:
        prompt = f"""
        Mavzu: {topic}
        Insho turi: {essay_type}
        Vazifa: O'zbek tilida chuqur mazmunli, badiiy va akademik uslubda insho yozing.
        
        Tuzilishi:
        1. Kirish (Tezis bayoni)
        2. Asosiy qism (Argumentlar va misollar)
        3. Xulosa
        
        Formatlash: Markdown formatidan foydalaning.
        """
        return await AIService.generate_content(prompt)

    @staticmethod
    async def generate_test(topic: str, count: str) -> str:
        prompt = f"""
        Mavzu: {topic}
        Savollar soni: {count} ta
        Vazifa: O'zbek tilida test savollari tuzing.
        
        Har bir savol uchun:
        - Savol matni
        - 4 ta variant (A, B, C, D)
        - To'g'ri javob
        - Qisqa izoh (nima uchun bu javob to'g'ri)
        
        Formatlash: Markdown formatida jadval yoki ro'yxat ko'rinishida.
        """
        return await AIService.generate_content(prompt)

# Bot Setup
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher(storage=MemoryStorage())
router = Router()

# Web App Setup
app = FastAPI(title="Talaba Bot API")

@app.get("/", response_class=HTMLResponse)
async def web_app_home():
    try:
        with open('webapp.html', 'r', encoding='utf-8') as f:
            return f.read()
    except FileNotFoundError:
        return "WebApp HTML fayli topilmadi."

@app.post("/api/user-info")
async def get_user_info(request: Request):
    # In a real app, you would validate the Telegram WebApp Init Data here
    # For now, we simulate pulling the ID, assuming the front end sends it or we use a default for testing
    # Since headers are hard to spoof in pure browser without tools, but for production use `validate_web_app_data`
    
    # We'll try to extract from initData if passed, otherwise default for testing
    body = await request.json() if await request.body() else {}
    # NOTE: In production, rely on initData from headers/body and verify signature
    
    # For simulation/demo purposes, we use a fixed ID if not provided, or extract from request
    tg_id = 123456789 # Default test ID
    
    user = await db.get_user(tg_id)
    if not user:
        return {
            "tg_id": tg_id,
            "balance": 0,
            "is_premium": False,
            "premium_until": None,
            "content_count": 0
        }
        
    return {
        "tg_id": user['tg_id'],
        "balance": user['balance'],
        "is_premium": bool(user['is_premium']),
        "premium_until": user['premium_until'],
        "content_count": user['content_count']
    }

@app.post("/api/generate")
async def generate_content_api(request: Request):
    data = await request.json()
    content_type = data.get('type')
    topic = data.get('topic')
    params = data.get('params', {})
    tg_id = 123456789 # Placeholder, should be extracted from auth
    
    if content_type == 'referat':
        content = await AIService.generate_referat(topic, params.get('size', '10'))
    elif content_type == 'prezentatsiya':
        content = await AIService.generate_presentation(topic, params.get('slides', '10'))
    elif content_type == 'insho':
        content = await AIService.generate_essay(topic, params.get('type', 'argumentativ'))
    elif content_type == 'test':
        content = await AIService.generate_test(topic, params.get('count', '10'))
    else:
        return {"success": False, "message": "Noto'g'ri kontent turi"}
    
    # Save statistics
    await db.save_content(tg_id, content_type, topic, content)
    await db.increment_content_count(tg_id)
    
    return {
        "success": True,
        "content": content,
        "type": content_type
    }



from pptx import Presentation
from pptx.util import Inches, Pt
from fastapi.responses import StreamingResponse
import io
import re

# PPTX Generator Helper
def create_pptx_from_text(content: str) -> io.BytesIO:
    prs = Presentation()
    
    # Split content into slides based on markdown headings or explicit separators
    # We look for lines starting with "Slide", "Slayd" or just headings
    slides_content = re.split(r'(?i)(?:^|\n)(?:Slide|Slayd)\s+\d+', content)
    
    # If split didn't work well (no explicit "Slide X"), try splitting by H1/H2 (#)
    if len(slides_content) < 2:
        slides_content = re.split(r'(?i)(?:^|\n)#+\s+', content)

    # Clean up empty splits
    slides_content = [s.strip() for s in slides_content if s.strip()]

    # If still no splits, treat whole text as one slide (fallback)
    if not slides_content:
        slides_content = [content]

    for i, slide_text in enumerate(slides_content):
        # Create a slide
        # 0 = Title Slide (for first slide if it looks like a title), 1 = Title and Content
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        
        # Extract Title
        lines = slide_text.split('\n')
        title_text = lines[0].strip().replace('#', '').strip()
        if title_text.startswith("*") and title_text.endswith("*"): # Remove bold markdown
            title_text = title_text.strip("*")
            
        # Remove title from lines to process body
        body_lines = lines[1:]
        body_text = '\n'.join(body_lines).strip()
        
        # Set Title
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            
        # Set Body Content
        if i == 0: # Title Slide
            # For title slide, put the rest in subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = body_text
        else: # Content Slide
            # For content slides, we want structured bullet points
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = body_text # Simple approach: dump text.

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

@app.post("/api/download-pptx")
async def download_pptx_api(request: Request):
    data = await request.json()
    content = data.get('content')
    topic = data.get('topic', 'Presentation')
    
    if not content:
        return JSONResponse(status_code=400, content={"message": "Content is required"})
        
    try:
        pptx_file = create_pptx_from_text(content)
        
        # Clean filename
        filename = re.sub(r'[^\w\-_\. ]', '_', topic) + ".pptx"
        
        return StreamingResponse(
            pptx_file, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        logger.error(f"PPTX Gen Error: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

# Bot Handlers

def get_main_menu():
    builder = ReplyKeyboardBuilder()
    builder.row(types.KeyboardButton(text="🌐 Web App", web_app=WebAppInfo(url=WEBAPP_URL)))
    builder.row(types.KeyboardButton(text="💰 To'lovlar"), types.KeyboardButton(text="👊 Admin"))
    return builder.as_markup(resize_keyboard=True)

@router.message(CommandStart())
async def cmd_start(message: Message):
    await db.create_or_update_user(message.from_user.id, message.from_user.username, message.from_user.full_name)
    await message.answer(
        f"👋 Assalomu alaykum, {message.from_user.full_name}!\n\n"
        "🤖 **Talaba Servis Bot** ga xush kelibsiz.\n\n"
        "Men orqali siz:\n"
        "📝 Referatlar\n"
        "📊 Prezentatsiyalar\n"
        "✍️ Insholar\n"
        "🧪 Testlar yaratishingiz mumkin.\n\n"
        "Eng qulay usul - **Web App** tugmasini bosish!",
        reply_markup=get_main_menu()
    )

from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup

class PaymentState(StatesGroup):
    waiting_for_receipt = State()

@router.message(F.text == "💰 To'lovlar")
async def cmd_payments(message: Message):
    text = (
        "💎 **Talaba Premium — Imkoniyatlaringizni kengaytiring!**\n\n"
        "Quyidagi tariflardan birini tanlang:\n\n"
        "💥 **1 Oy** – 25 000 so‘m\n"
        "⚡ **3 Oy** – 60 000 so‘m 🔻 (20% Tejang!)\n"
        "🌟 **6 Oy** – 99 000 so‘m 🔻 (34% Tejang!)\n\n"
        "Tanlash uchun tugmalardan foydalaning:"
    )
    
    builder = InlineKeyboardBuilder()
    builder.button(text="💥 1 Oy (25k)", callback_data="pay_1m")
    builder.button(text="⚡ 3 Oy (60k)", callback_data="pay_3m")
    builder.button(text="🌟 6 Oy (99k)", callback_data="pay_6m")
    builder.adjust(2, 1)
    
    await message.answer(text, reply_markup=builder.as_markup())

@router.callback_query(F.data.in_({"pay_1m", "pay_3m", "pay_6m"}))
async def process_tariff_selection(callback: CallbackQuery, state: FSMContext):
    plan_map = {
        "pay_1m": ("1 Oylik", 25000),
        "pay_3m": ("3 Oylik", 60000),
        "pay_6m": ("6 Oylik", 99000)
    }
    
    plan, amount = plan_map[callback.data]
    await state.update_data(plan=plan, amount=amount)
    await state.set_state(PaymentState.waiting_for_receipt)
    
    text = (
        "💳 **To‘lov ma’lumotlari**\n\n"
        f"🏷 **Tarif:** {plan}\n"
        f"💵 **Summa:** {amount:,} so‘m\n"
        "💳 **Karta:** (U.X. nomida)\n\n"
        "➡ `5614 6821 1601 2602` (Uzcard)\n"
        "➡ `4023 0605 1561 2637` (Visa)\n\n"
        "➤ To‘lovni amalga oshirgach, chekni (rasm yoki fayl) shu chatga yuboring. "
        "Chek tasdiqlash uchun botga yuboriladi.\n\n"
        "❌ Bekor qilish uchun “Bekor qilish” tugmasini bosing."
    ).replace(",", " ")
    
    builder = InlineKeyboardBuilder()
    builder.button(text="❌ Bekor qilish", callback_data="cancel_payment")
    
    # Edit the message or send new one? Ideally send new one or edit.
    # Image 2 shows a new message potentially, but inline editing is cleaner.
    # Let's delete the old one and send a new one to mimic the screenshot flow perfectly, 
    # OR edit. Editing is usually better UX in bots.
    await callback.message.edit_text(text, reply_markup=builder.as_markup())
    await callback.answer()

@router.callback_query(F.data == "cancel_payment")
async def cancel_payment_handler(callback: CallbackQuery, state: FSMContext):
    await state.clear()
    await callback.message.delete()
    await callback.answer("To'lov bekor qilindi")

@router.message(F.photo, PaymentState.waiting_for_receipt)
async def handle_payment_receipt(message: Message, state: FSMContext):
    """Handle photo receipts sent by user with state."""
    data = await state.get_data()
    amount = data.get("amount", 50000)
    plan = data.get("plan", "Noma'lum")
    
    photo = message.photo[-1]
    file_id = photo.file_id
    payment_id = str(uuid.uuid4())
    
    # Create payment record
    await db.create_payment(
        tg_id=message.from_user.id,
        amount=amount,
        card_number=f"PHOTO_RECEIPT_{plan}",
        card_holder=message.from_user.full_name,
        payment_id=payment_id
    )
    
    # Update admin_note with file_id
    async with aiosqlite.connect('talaba_bot.db') as conn:
        await conn.execute("UPDATE payments SET admin_note=? WHERE payment_id=?", (file_id, payment_id))
        await conn.commit()
    
    await state.clear()
    await message.answer("✅ To'lov cheki qabul qilindi! Admin tasdiqlashini kuting.")
    
    try:
        builder = InlineKeyboardBuilder()
        builder.button(text="✅ Tasdiqlash", callback_data=f"approve_{payment_id}")
        builder.button(text="❌ Rad etish", callback_data=f"reject_{payment_id}")
        
        caption = (
            f"💰 **Yangi To'lov Cheki!**\n\n"
            f"👤: {message.from_user.full_name}\n"
            f"🆔: `{message.from_user.id}`\n"
            f"🏷: {plan}\n"
            f"💵: {amount:,} so'm"
        ).replace(",", " ")

        await bot.send_photo(
            ADMIN_ID,
            photo=file_id,
            caption=caption,
            reply_markup=builder.as_markup()
        )
    except Exception as e:
        logger.error(f"Failed to notify admin: {e}")
    



@router.message(F.text == "👊 Admin")
async def cmd_admin(message: Message):
    if message.from_user.id != ADMIN_ID:
        await message.answer("⛔ Siz admin emassiz.")
        return
        
    stats = await db.get_stats()
    text = f"""
📊 **Admin Statistika**

👥 Foydalanuvchilar: {stats['users']}
📝 Yaratilgan kontent: {stats['content']}
💰 Kutilayotgan to'lovlar: {stats['pending_payments']}
    """
    
    keyboard = InlineKeyboardBuilder()
    if stats['pending_payments'] > 0:
        keyboard.button(text="To'lovlarni ko'rish", callback_data="view_payments")
    
    await message.answer(text, reply_markup=keyboard.as_markup())

@router.callback_query(F.data == "view_payments")
async def view_payments_callback(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
        
    payments = await db.get_pending_payments()
    if not payments:
        await callback.message.answer("Kutilayotgan to'lovlar qolmadi.")
        return
        
    for p in payments:
        builder = InlineKeyboardBuilder()
        builder.button(text="✅ Tasdiqlash", callback_data=f"approve_{p['payment_id']}")
        builder.button(text="❌ Rad etish", callback_data=f"reject_{p['payment_id']}")
        builder.adjust(2)
        
        file_id = p['admin_note']
        caption = (f"👤 Foydalanuvchi: {p['card_holder']}\n"
                   f"🆔 ID: {p['tg_id']}\n"
                   f"💰 Summa: {p['amount']} so'm\n"
                   f"📅 Vaqt: {p['created']}")
        
        if file_id and len(file_id) > 10:
            await callback.message.answer_photo(
                photo=file_id,
                caption=caption,
                reply_markup=builder.as_markup()
            )
        else:
            await callback.message.answer(
                text=caption,
                reply_markup=builder.as_markup()
            )
    await callback.answer()

@router.callback_query(F.data.startswith("approve_"))
async def approve_payment_handler(callback: CallbackQuery):
    payment_id = callback.data.split("_")[1]
    payment = await db.get_payment(payment_id)
    
    if payment:
        await db.update_payment_status(payment_id, "approved")
        await db.set_premium(payment['tg_id'])
        
        try:
            await bot.send_message(payment['tg_id'], "✅ Sizning to'lovingiz tasdiqlandi! Premium status faollashtirildi.")
        except:
            pass
            
        new_text = f"✅ To'lov tasdiqlandi ({payment['amount']} so'm)\n👤 {payment['card_holder']}"
        if callback.message.photo:
            await callback.message.edit_caption(caption=new_text, reply_markup=None)
        else:
            await callback.message.edit_text(text=new_text, reply_markup=None)
    else:
        await callback.message.answer("❌ To'lov topilmadi")

@router.callback_query(F.data.startswith("reject_"))
async def reject_payment_handler(callback: CallbackQuery):
    payment_id = callback.data.split("_")[1]
    payment = await db.get_payment(payment_id)
    
    if payment:
        await db.update_payment_status(payment_id, "rejected")
        try:
            await bot.send_message(payment['tg_id'], "❌ Sizning to'lovingiz rad etildi.")
        except:
            pass
            
        new_text = f"❌ To'lov rad etildi\n👤 {payment['card_holder']}"
        if callback.message.photo:
            await callback.message.edit_caption(caption=new_text, reply_markup=None)
        else:
            await callback.message.edit_text(text=new_text, reply_markup=None)
    else:
        await callback.message.answer("❌ To'lov topilmadi")

# Main Function
async def main():
    # Initialize DB
    await db.init_db()
    
    # Include Router
    dp.include_router(router)
    
    # Start Web App Server in Background
    # NOTE: In production, run uvicorn separately. Here we run it in a task for the bot process.
    config = uvicorn.Config(app, host="0.0.0.0", port=8081, log_level="info")
    server = uvicorn.Server(config)
    
    # Run server and bot concurrently
    await asyncio.gather(
        server.serve(),
        dp.start_polling(bot)
    )

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        logger.info("Bot stopped")
    except Exception as e:
        logger.error(f"Error: {e}")
