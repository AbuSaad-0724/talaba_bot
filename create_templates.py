#!/usr/bin/env python3
"""
Template Creator
Creates sample PowerPoint templates for the bot
"""

from utils.template_gen import create_starter_template
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os

def create_templates():
    """Create multiple template variations"""
    print("🎨 Creating PowerPoint Templates...\n")
    
    os.makedirs("data/templates", exist_ok=True)
    
    # Template 1: Default (already exists)
    print("Template 1: Default Template")
    print("-" * 60)
    try:
        create_starter_template()
        if os.path.exists("data/templates/default.pptx"):
            size = os.path.getsize("data/templates/default.pptx")
            print(f"✅ Created: data/templates/default.pptx ({size:,} bytes)\n")
        else:
            print("❌ Failed to create default template\n")
    except Exception as e:
        print(f"❌ Error: {e}\n")
    
    # Template 2: Academic Style
    print("Template 2: Academic Style")
    print("-" * 60)
    try:
        prs = Presentation()
        
        # Title slide
        title_layout = prs.slide_master.slide_layouts[0]
        bg = title_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(25, 25, 112)  # Midnight Blue
        
        # Content slide
        content_layout = prs.slide_master.slide_layouts[1]
        bg = content_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # White Smoke
        
        output = "data/templates/academic.pptx"
        prs.save(output)
        size = os.path.getsize(output)
        print(f"✅ Created: {output} ({size:,} bytes)\n")
    except Exception as e:
        print(f"❌ Error: {e}\n")
    
    # Template 3: Modern Style
    print("Template 3: Modern Style")
    print("-" * 60)
    try:
        prs = Presentation()
        
        # Title slide - Gradient effect simulation
        title_layout = prs.slide_master.slide_layouts[0]
        bg = title_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(70, 130, 180)  # Steel Blue
        
        # Content slide
        content_layout = prs.slide_master.slide_layouts[1]
        bg = content_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice Blue
        
        output = "data/templates/modern.pptx"
        prs.save(output)
        size = os.path.getsize(output)
        print(f"✅ Created: {output} ({size:,} bytes)\n")
    except Exception as e:
        print(f"❌ Error: {e}\n")
    
    # Template 4: Dark Theme
    print("Template 4: Dark Theme")
    print("-" * 60)
    try:
        prs = Presentation()
        
        # Title slide
        title_layout = prs.slide_master.slide_layouts[0]
        bg = title_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark Gray
        
        # Content slide
        content_layout = prs.slide_master.slide_layouts[1]
        bg = content_layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(45, 45, 45)  # Darker Gray
        
        output = "data/templates/dark.pptx"
        prs.save(output)
        size = os.path.getsize(output)
        print(f"✅ Created: {output} ({size:,} bytes)\n")
    except Exception as e:
        print(f"❌ Error: {e}\n")
    
    # Summary
    print("=" * 60)
    print("✅ Template creation completed!")
    print("\nAvailable templates:")
    
    template_dir = "data/templates"
    if os.path.exists(template_dir):
        for file in os.listdir(template_dir):
            if file.endswith('.pptx'):
                path = os.path.join(template_dir, file)
                size = os.path.getsize(path)
                print(f"  • {file}: {size:,} bytes")

if __name__ == "__main__":
    create_templates()
