from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

TEMPLATE_PATH = "data/templates/default.pptx"

def create_starter_template():
    """Create a stylish starter template if it doesn't exist."""
    os.makedirs("data/templates", exist_ok=True)
    
    prs = Presentation()
    
    # Set slide width and height for standard 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. Title Slide Layout
    if len(prs.slide_master.slide_layouts) > 0:
        slide_layout = prs.slide_master.slide_layouts[0]
        try:
            background = slide_layout.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue background
        except Exception:
            pass  # Some templates don't allow background modification
        
        # Customize Title Text Style
        try:
            title_style = slide_layout.placeholders[0].text_frame.paragraphs[0].font
            title_style.name = "Arial"
            title_style.size = Pt(44)
            title_style.color.rgb = RGBColor(255, 255, 255) # White text
            
            # Customize Subtitle
            subtitle = slide_layout.placeholders[1].text_frame.paragraphs[0].font
            subtitle.name = "Arial"
            subtitle.color.rgb = RGBColor(200, 200, 200) # Light gray
        except Exception:
            pass  # Placeholder indices may vary
    
    # 2. Content Slide Layout
    if len(prs.slide_master.slide_layouts) > 1:
        bullet_layout = prs.slide_master.slide_layouts[1]
        try:
            background = bullet_layout.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(240, 248, 255) # AliceBlue (Very light)
        except Exception:
            pass
        
        # Title on content slide
        try:
            title_c = bullet_layout.placeholders[0].text_frame.paragraphs[0].font
            title_c.name = "Arial"
            title_c.color.rgb = RGBColor(0, 51, 102) # Dark Blue
            title_c.bold = True
        except Exception:
            pass
    
    try:
        prs.save(TEMPLATE_PATH)
        print(f"Starter template created at {TEMPLATE_PATH}")
    except Exception as e:
        print(f"Error saving template: {e}")
        # Try saving to a different location
        fallback_path = "default.pptx"
        prs.save(fallback_path)
        print(f"Fallback template saved to {fallback_path}")

if __name__ == "__main__":
    create_starter_template()
