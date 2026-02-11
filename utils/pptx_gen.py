from pptx import Presentation
from pptx.util import Inches, Pt
import os
import re
import logging
from utils.template_gen import create_starter_template, TEMPLATE_PATH

logger = logging.getLogger(__name__)

def create_presentation_pptx(topic: str, content: str, output_path: str, template_path: str = None):
    if not template_path:
        if not os.path.exists(TEMPLATE_PATH):
            try:
                create_starter_template()
            except Exception as e:
                logger.error(f"Failed to create template: {e}")
                template_path = None
        
        if not template_path or not os.path.exists(TEMPLATE_PATH if template_path else TEMPLATE_PATH):
            template_path = None
    
    try:
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            logger.warning("Using blank presentation - no template found")
    except Exception as e:
        logger.error(f"PPTX Template load error: {e}. Using blank presentation.")
        prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = topic
    subtitle.text = "Talaba Bot tomonidan tayyorlandi"
    
    # Content Slides
    # Strategy 1: Try strict delimiter "|||"
    sections = content.split("|||")
    
    # Strategy 2: If strict fail, try "Slayd X:" pattern
    if len(sections) < 2:
        logger.warning("PPTX Gen: '|||' delimiter failed. Trying 'Slayd' pattern.")
        sections = re.split(r'(?:^|\n)(?:Slayd|Slide)\s*\d+[:\.]?', content, flags=re.IGNORECASE)

    # Strategy 3: Just split by double newlines if all else fails
    if len(sections) < 2:
        logger.warning("PPTX Gen: 'Slayd' pattern failed. Trying double newlines.")
        sections = content.split("\n\n")

    valid_slides_count = 0
    for section in sections:
        section = section.strip()
        if not section: continue
        
        # Skip garbage like just numbers or "Slayd 1" headers without content
        if len(section) < 5: continue
        
        lines = section.split("\n", 1)
        slide_title = lines[0].strip()
        slide_text = lines[1].strip() if len(lines) > 1 else ""
        
        # Cleanup title
        # Remove common prefixes like "1.", "1:", "Slayd 1:", etc.
        slide_title = re.sub(r'^(?:Slayd|Slide|Step|Bo\'lim)\s*\d*[:\.]?\s*', '', slide_title, flags=re.IGNORECASE).strip()
        slide_title = re.sub(r'^\d+[:\.]\s*', '', slide_title).strip()
        # Remove markdown bold/italic
        slide_title = slide_title.replace('**', '').replace('__', '').replace('*', '')
        
        # If text is empty, maybe the title line was too long and actually contains the text?
        if not slide_text and len(slide_title) > 50:
             # Basic heuristic: split by first punctuation
             parts = re.split(r'[:.?!]\s', slide_title, 1)
             if len(parts) > 1:
                 slide_title = parts[0]
                 slide_text = parts[1]

        # Check if slide layouts exist
        if len(prs.slide_layouts) < 2:
            logger.error("Not enough slide layouts in presentation")
            break
            
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_title
        tf = body_shape.text_frame
        tf.text = slide_text
        valid_slides_count += 1
        
        # Minimal styling
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            
    if valid_slides_count == 0:
        # Emergency fallback: Create one slide with everything
        logger.warning("No valid slides created, using emergency fallback")
        if len(prs.slide_layouts) > 1:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Taqdimot"
            slide.shapes.placeholders[1].text_frame.text = content[:1000] # Limit char count
        else:
            # Last resort - add to title slide
            slide.shapes.title.text = topic
            slide.placeholders[1].text_frame.text = content[:500]

    prs.save(output_path)
    return output_path
